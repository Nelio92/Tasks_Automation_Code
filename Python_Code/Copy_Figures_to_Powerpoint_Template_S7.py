# -*- coding: utf-8 -*-
"""
Script to import all figures from a specified folder to a Powerpoint template.
@Created on Wed Jul 25 14:53:26 2018
@Author: Wandji, Lionel Wilfried.
"""

"""     Powerpoint slides layouts """
# 0 Title (presentation title slide)
# 1 Title and Content
# 2 Section Header (sometimes called Segue)
# 3 Two Content (side by side bullet textboxes)
# 4 Comparison (same but additional title for each side by side content box)
# 5 Title Only
# 6 Blank
# 7 Content with Caption
# 8 Picture with Caption

"""     Import the right libraries """
from pptx import Presentation 
from pptx.util import Inches
import glob as gb
import os
import pandas as pd
import numpy as np
import shutil
from os.path import basename

"""     Folder Path """
print("")
print("Start...")
#img_path = 'C:/UserData/WinPython/testPics/*.png' # The folder containing all your plots (CDF, Boxes and PCsites) arranged by Name
#img_path2 = 'C:/UserData/WinPython/testPics/'
#img_path = 'C:/UserData/WinPython/Pics_S7/*.png' # The folder containing all your plots (CDF, Boxes and PCsites) arranged by Name
#img_path2 = 'C:/UserData/WinPython/Pics_S7/'
#img_sorted_path = 'C:/UserData/WinPython/Pics_Sorted_S7/'
img_path = 'C:/UserData/WinPython/Pics_S7_SHORT/*.png' # The folder containing all your plots (CDF, Boxes and PCsites) arranged by Name
img_path2 = 'C:/UserData/WinPython/Pics_S7_SHORT/'
img_sorted_path = 'C:/UserData/WinPython/Pics_Sorted_S7_SHORT/'
files = gb.glob(img_path)
# Sort the file names alphanumerically from the smallest to the highest TEST NUMBERS in new directory
files2 = []
for file in files:
    fileNameWithExtension = os.path.splitext(basename(file))
    fileName = fileNameWithExtension[0]
    oldDir = os.path.join(img_path2, fileName + '.png')
    if (fileName[0:5].isdigit() == False):
        fileName2 = '0' + fileName
        if os.path.exists(os.path.join(img_sorted_path, fileName2 + '.png')) == False:
            shutil.copy(oldDir, os.path.join(img_sorted_path, fileName + '.png'))
            os.rename(os.path.join(img_sorted_path, fileName + '.png'), os.path.join(img_sorted_path, fileName2 + '.png'))
            files2.append(os.path.join(img_sorted_path, fileName2 + '.png'))
        else:
            files2.append(os.path.join(img_sorted_path, fileName2 + '.png'))
    else:
        if os.path.exists(os.path.join(img_sorted_path, fileName + '.png')) == False:
            sortedDir = os.path.join(img_sorted_path, fileName + '.png')
            shutil.copy(oldDir, os.path.join(img_sorted_path, fileName + '.png'))
            files2.append(sortedDir)
        else:
            files2.append(os.path.join(img_sorted_path, fileName + '.png'))
files2 = sorted(files2)
print("")
print("Copying and sorting of files finished...")

"""     Create a new presentation, import figures and apply formatting """

prs = Presentation('template_S7.pptx') # I took the template of Osama with just the COVER and BACK pages
slideLayout = prs.slide_layouts[6]
counter = 0
# Check the yield losses in Excel file
YL_CSV = pd.read_excel('yield_losses_all_tests_S7.xlsx', sheet_name="yield_losses_all_tests")
df = pd.DataFrame(YL_CSV)
YL_TestNumber = np.array(df.iloc[:, 0].values)
YL_TestNames = np.array(df.iloc[:, 1].values)
YL_PATType = np.array(df.iloc[:, 3].values)
YL = np.array(df.iloc[:, 4].values)
Cpk_DPAT = np.array(df.iloc[:, 5].values)
Cpk_SPAT = np.array(df.iloc[:, 6].values)
LSL = np.array(df.iloc[:, 7].values)
USL = np.array(df.iloc[:, 8].values)
SPAT_LL = np.array(df.iloc[:, 9].values)
SPAT_UL = np.array(df.iloc[:, 10].values)
Distribution = np.array(df.iloc[:, 11].values)

Passed_Devices = 568068

# Slides
for file in files2:
    fileName = os.path.splitext(basename(file))
    if "3A8" in fileName[0]:
        continue
    elif ("S7 box lot.png" in file) or ("S7 box lot zoom.png" in file) or \
          ("S7 CDF lot.png" in file) or ("S7 CDF lot zoom.png" in file) or \
          ("S7 CDF site.png" in file) or ("S7 CDF site zoom.png" in file): 
        for num, nam, pat, yl, cpkD, cpkS, lsl, usl, sll, sul, dis in zip(YL_TestNumber, YL_TestNames, YL_PATType, YL, Cpk_DPAT, Cpk_SPAT, LSL, USL, SPAT_LL, SPAT_UL, Distribution):
            lsName = str(num) + ', ' + nam
            if (lsName in fileName[0]) and (pat == 'DPAT'):
#                if  yl == 0:
#                    continue
#                else:
                    # Inserting figures here on blank slides
                    slide = prs.slides.add_slide(slideLayout)
                    titleForSlide = fileName[0]
                    pic = slide.shapes.add_picture(file, top=Inches(1), left=Inches(0.2), height=Inches(6))
                    title = slide.shapes.title
                    title.text = titleForSlide
                    counter += 1
                    if ("S7 CDF lot.png" in file):
                        # Small table for the Yield Loss
                        table = slide.shapes.add_table(10, 2, left=Inches(6.2), top=Inches(1.3), width=Inches(3.5), height=Inches(2)).table
                        table.cell(0, 0).text = 'DPAT YL (%)' # Yield Loss as YL
                        table.cell(0, 1).text = str(round(yl*100, 5))
                        table.cell(1, 0).text = '#PAT_Outliers' # Yield Loss as YL
                        table.cell(1, 1).text = str(round(yl*Passed_Devices))
                        table.cell(2, 0).text = 'Cpk SPAT' # Yield Loss as YL
                        table.cell(2, 1).text = str(round(cpkS, 2))
                        table.cell(3, 0).text = 'Cpk DPAT' # Yield Loss as YL
                        table.cell(3, 1).text = str(round(cpkD, 2))
                        table.cell(4, 0).text = 'LSL' # Yield Loss as YL
                        table.cell(4, 1).text = str(lsl)
                        table.cell(5, 0).text = 'USL' # Yield Loss as YL
                        table.cell(5, 1).text = str(usl)
                        table.cell(6, 0).text = 'SPAT LL' # Yield Loss as YL
                        table.cell(6, 1).text = str(round(sll, 4))
                        table.cell(7, 0).text = 'SPAT UL' # Yield Loss as YL
                        table.cell(7, 1).text = str(round(sul, 4))
                        table.cell(8, 0).text = 'Category' # Yield Loss as YL
                        table.cell(8, 1).text = str(dis)
                        table.cell(9, 0).text = 'Remarks' # Yield Loss as YL
                        table.cell(9, 1).text = ''                       
    else:
        continue
print("")    
print("{} figures imported to the presentation.".format(counter))

"""     Save the presentation """
prs.save('Review_PAT_Analysis_W1390C_S7.pptx') # give the name you want to your specified folder 
print("")
print("End...")