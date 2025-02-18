# -*- coding: utf-8 -*-
"""
Script to import all figures from a specified folder to a Powerpoint presentation.
Author: Wandji, Lionel Wilfried.
Date:   07.2018 
"""

"""     Powerpoint slides layouts """
# 0 Title (presentation title slide)
# 1 Title and Content
# 2 Section Header (sometimes called Segue)
# 3 Two Content (side by side bullet textboxes)
# 4 Comparison (same but additional title for each side by side content box)
# 5 Title Only
# 6 Blank
# 7 Content with Caption
# 8 Picture with Caption

"""     Import the right libraries """
from pptx import Presentation 
from pptx.util import Inches
import glob as gb
import os
from os.path import basename

"""     Folder Path """
img_path = 'C:/UserData/WinPython/Pics/*.png'
files = gb.glob(img_path)

"""     Create a new presentation, import figures and apply formatting """
prs = Presentation('template.pptx')
slideLayout = prs.slide_layouts[6]
counter = 0
# Slides
for file, picNr in zip(files, range(1, len(files)+1)):
    fileName = os.path.splitext(basename(file))
    if "3A8" in fileName[0]:
        continue
    elif ("S6 box lot zoom.png" in file) or ("S6 box lot.png" in file) or \
          ("S6 CDF lot zoom.png" in file) or ("S6 CDF lot.png" in file) or \
          ("S6 CDF site zoom.png" in file) or ("S6 CDF site.png" in file):
              
        # Inserting figures here on blank slides
        slide = prs.slides.add_slide(slideLayout)
        titleForSlide = fileName[0]
        pic = slide.shapes.add_picture(file, top=Inches(1), left=Inches(0.5), height=Inches(6))
        title = slide.shapes.title
        title.text = titleForSlide
        print(titleForSlide)
        counter += 1
        if ("S6 box lot zoom.png" in file):
            # Small table for the Yield Loss
            table = slide.shapes.add_table(2, 1, left=Inches(9), top=Inches(2), width=Inches(0.75), height=Inches(1)).table
            #table.columns[0].width = Inches(1)
            #table.columns[1].width = Inches(1)
            table.cell(0, 0).text = 'YL'
            table.cell(1, 0).text = ''
    else:
        continue
print("")    
print("{} figures imported to the presentation.".format(counter))

"""     Save the presentation """
prs.save('Presentation_Review_PAT_Analysis_W13__S.pptx')

