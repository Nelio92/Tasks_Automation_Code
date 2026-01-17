from __future__ import annotations

import math
from dataclasses import dataclass
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[2]
REPORT_DIR = Path(__file__).resolve().parent
ASSETS_DIR = REPORT_DIR / "assets"


@dataclass(frozen=True)
class ReportConfig:
    title: str = "ATE ↔ CV Correlation Automation Report"
    subtitle: str = "Procedure, deliverables, and time benefit (generated from Python automation scripts)"
    author: str = "Auto-generated"

    # Assumptions for time estimate (explicitly shown in the report).
    manual_minutes_per_iteration: float = 180.0
    automated_minutes_per_iteration: float = 25.0
    iterations_per_week: int = 2


def _pick_first(glob_pattern: str) -> Path | None:
    hits = sorted(REPO_ROOT.glob(glob_pattern))
    for p in hits:
        if p.is_file():
            return p
    return None


def _save_flow_diagram_png(path: Path) -> None:
    fig = plt.figure(figsize=(14, 4.2), dpi=160)
    ax = fig.add_subplot(111)
    ax.axis("off")

    boxes = [
        (0.04, 0.55, 0.18, 0.30, "Raw ATE exports\n(.xlsx / .csv)\n+ chip list"),
        (0.28, 0.55, 0.18, 0.30, "1) Tests_Data_Extractor_Flat\nFilters chips + tests\nAdds meta (Temp/Vcorner/Freq)"),
        (0.52, 0.55, 0.18, 0.30, "2) Correlation scripts\nDPLL: delta-based limits\nModelX: Offset + Physics(Kf)"),
        (0.76, 0.55, 0.20, 0.30, "Deliverables\nExcel summaries\n+ per-group plots"),
    ]

    for x, y, w, h, label in boxes:
        rect = plt.Rectangle((x, y), w, h, fill=False, linewidth=2)
        ax.add_patch(rect)
        ax.text(x + w / 2, y + h / 2, label, ha="center", va="center", fontsize=11)

    def arrow(x1, y1, x2, y2):
        ax.annotate(
            "",
            xy=(x2, y2),
            xytext=(x1, y1),
            arrowprops=dict(arrowstyle="->", linewidth=2),
        )

    arrow(0.22, 0.70, 0.28, 0.70)
    arrow(0.46, 0.70, 0.52, 0.70)
    arrow(0.70, 0.70, 0.76, 0.70)

    ax.text(
        0.5,
        0.18,
        "Key idea: replace repeated manual Excel filtering/merging/plotting with deterministic scripts\n"
        "→ consistent grouping, traceable outputs, and easy reruns for new lots / DoE splits.",
        ha="center",
        va="center",
        fontsize=11,
    )

    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def _save_time_benefit_png(path: Path, cfg: ReportConfig) -> None:
    manual_h = cfg.manual_minutes_per_iteration / 60.0
    auto_h = cfg.automated_minutes_per_iteration / 60.0
    weekly_manual = manual_h * cfg.iterations_per_week
    weekly_auto = auto_h * cfg.iterations_per_week
    weekly_saved = weekly_manual - weekly_auto

    fig = plt.figure(figsize=(10, 4.8), dpi=160)
    ax = fig.add_subplot(111)

    labels = ["Manual", "Automated"]
    hours = [weekly_manual, weekly_auto]

    ax.bar(labels, hours)
    ax.set_ylabel("Hours per week")
    ax.set_title("Estimated time per week (assumptions stated on slide)")
    ax.grid(True, axis="y", alpha=0.25)

    for i, v in enumerate(hours):
        ax.text(i, v + 0.05, f"{v:.2f}h", ha="center", va="bottom", fontsize=11)

    ax.text(
        0.5,
        -0.22,
        f"Estimated weekly time saved: {weekly_saved:.2f}h/week",
        transform=ax.transAxes,
        ha="center",
        va="top",
        fontsize=12,
    )

    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def _add_title_slide(prs: Presentation, cfg: ReportConfig) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = cfg.title
    subtitle = slide.placeholders[1]
    subtitle.text = cfg.subtitle


def _add_picture_slide(prs: Presentation, title: str, image_path: Path, *, caption: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    slide.shapes.title.text = title

    # image sized for 16:9 (13.33" x 7.5")
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(12.33)
    slide.shapes.add_picture(str(image_path), left, top, width=width)

    if caption:
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.33), Inches(0.4)).text_frame
        p = tx.paragraphs[0]
        p.text = caption
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER


def _add_two_pictures_slide(
    prs: Presentation,
    title: str,
    left_img: Path,
    right_img: Path,
    *,
    left_caption: str = "",
    right_caption: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    top = Inches(1.2)
    img_w = Inches(6.15)
    slide.shapes.add_picture(str(left_img), Inches(0.5), top, width=img_w)
    slide.shapes.add_picture(str(right_img), Inches(6.68), top, width=img_w)

    if left_caption:
        tf = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), img_w, Inches(0.4)).text_frame
        tf.text = left_caption
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    if right_caption:
        tf = slide.shapes.add_textbox(Inches(6.68), Inches(7.0), img_w, Inches(0.4)).text_frame
        tf.text = right_caption
        tf.paragraphs[0].font.size = Pt(11)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER


def _add_bullets_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, b in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(18)


def _add_process_summary_slide(prs: Presentation) -> None:
    # Keep this mostly visual: a compact, structured bullet list.
    _add_bullets_slide(
        prs,
        "Correlation procedure (what the scripts actually do)",
        [
            "Extract ATE raw exports → normalize → filter to target chips/tests",
            "Enrich metadata: Temperature from filename; Voltage corner + Frequency from Test Name",
            "Merge CV/ATE rows by keys (Wafer/X/Y/DUT/Temp/Vcorner/Freq/Test#)",
            "Group by test-case (TXLO: Test#…; TXPA: LUT value… + optional PA channel for LUT255)",
            "Compute correlation + new limits and output: Excel summaries + per-group plots",
        ],
    )


def _add_automation_highlights_slide(prs: Presentation) -> None:
    _add_bullets_slide(
        prs,
        "Automation highlights vs manual workflow",
        [
            "Robust chip-list ingestion (CSV/XLSX; flexible headers; DoE split + DUT Nr preserved)",
            "Test selection supports explicit IDs, ranges, and substrings",
            "Batch-processing of all input files + automatic skip of temp/~$ files",
            "Consistent, reproducible grouping and output naming (traceable plots per test-case)",
            "Limit computation and guardbanding performed systematically (no copy/paste errors)",
        ],
    )


def _add_time_slide(prs: Presentation, cfg: ReportConfig, time_img: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Time benefit estimate (automation impact)"

    slide.shapes.add_picture(str(time_img), Inches(0.7), Inches(1.25), width=Inches(7.2))

    # Assumptions box (small, explicit)
    tf = slide.shapes.add_textbox(Inches(8.1), Inches(1.35), Inches(4.7), Inches(5.8)).text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "Assumptions (edit to match reality):"
    p0.font.bold = True
    p0.font.size = Pt(16)

    assumptions = [
        f"Manual effort per iteration: {cfg.manual_minutes_per_iteration:.0f} min",
        f"Automated effort per iteration: {cfg.automated_minutes_per_iteration:.0f} min",
        f"Iterations per week: {cfg.iterations_per_week}",
        "Manual includes: Excel filtering/merging, plotting, copying results",
        "Automated includes: editing config + running scripts + quick review",
    ]

    for a in assumptions:
        p = tf.add_paragraph()
        p.text = a
        p.level = 0
        p.font.size = Pt(14)


def _add_outlook_slide(prs: Presentation) -> None:
    _add_bullets_slide(
        prs,
        "Outlook: suggested next automation steps",
        [
            "Central config file (YAML/JSON) + CLI wrapper to remove in-code edits",
            "Single pipeline runner: Extract → Correlate → Package plots + summary",
            "Add validation: missing columns, Kf merge coverage thresholds, inverted limit windows",
            "Standardize outputs: versioned run folders + metadata (git hash, date, dataset tags)",
            "Regression tests on small sample datasets to lock behavior",
        ],
    )


def main() -> Path:
    cfg = ReportConfig()
    ASSETS_DIR.mkdir(parents=True, exist_ok=True)

    # Pick representative real plots from existing output folders
    lo_series = _pick_first("plots_LO_Power_DoE/**/**__series.png")
    lo_models = _pick_first("plots_LO_Power_DoE/**/**__models.png")
    pa255_series = _pick_first("plots_PA_Power_DoE/**/LUT_value_255*__series.png")
    pa255_models = _pick_first("plots_PA_Power_DoE/**/LUT_value_255*__models.png")

    flow_png = ASSETS_DIR / "workflow.png"
    time_png = ASSETS_DIR / "time_benefit.png"
    _save_flow_diagram_png(flow_png)
    _save_time_benefit_png(time_png, cfg)

    prs = Presentation()
    # Set 16:9 wide
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, cfg)
    _add_picture_slide(
        prs,
        "End-to-end automation pipeline",
        flow_png,
        caption="From raw ATE exports to correlated limits and plots",
    )
    _add_process_summary_slide(prs)
    _add_automation_highlights_slide(prs)

    # Example outputs (real images if present)
    if lo_series and lo_models:
        _add_two_pictures_slide(
            prs,
            "Example output: TXLO (series + models)",
            lo_series,
            lo_models,
            left_caption="Series: CV vs ATE + correlated + limits",
            right_caption="Models: regression view + residuals",
        )

    if pa255_series and pa255_models:
        _add_two_pictures_slide(
            prs,
            "Example output: TXPA LUT255 (with requirements guardbanding)",
            pa255_series,
            pa255_models,
            left_caption="Series: correlated distribution + requirement lines",
            right_caption="Models: regression + residuals (offset/physics)",
        )

    _add_time_slide(prs, cfg, time_png)
    _add_outlook_slide(prs)

    out_path = REPORT_DIR / "ATE_CV_Correlation_Automation_Report.pptx"
    prs.save(out_path)
    return out_path


if __name__ == "__main__":
    p = main()
    print(f"Wrote: {p}")
